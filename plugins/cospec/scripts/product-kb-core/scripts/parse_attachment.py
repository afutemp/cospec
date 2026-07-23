#!/usr/bin/env python3
"""解析产品知识库采集到的常见文档附件。"""

import json
import os
import sys


def normalize_text(value):
    return "\n".join(
        line.rstrip()
        for line in value.replace("\r\n", "\n").replace("\r", "\n").split("\n")
    ).strip()


def parse_pdf(file_path):
    from pypdf import PdfReader

    reader = PdfReader(file_path)
    return "\n\n".join(page.extract_text() or "" for page in reader.pages)


def parse_docx(file_path):
    from docx import Document

    document = Document(file_path)
    lines = [paragraph.text for paragraph in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            lines.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(lines)


def parse_xlsx(file_path):
    from openpyxl import load_workbook

    workbook = load_workbook(file_path, read_only=True, data_only=True)
    sections = []
    try:
        for worksheet in workbook.worksheets:
            sections.append(f"## {worksheet.title}")
            for row in worksheet.iter_rows(values_only=True):
                values = ["" if value is None else str(value) for value in row]
                if any(values):
                    sections.append(" | ".join(values))
    finally:
        workbook.close()
    return "\n".join(sections)


def parse_pptx(file_path):
    from pptx import Presentation

    presentation = Presentation(file_path)
    lines = []
    for index, slide in enumerate(presentation.slides, start=1):
        lines.append(f"## Slide {index}")
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                lines.append(text)
    return "\n".join(lines)


def parse_document(file_path, extension_override=None):
    extension = (extension_override or os.path.splitext(file_path)[1]).lower()
    parsers = {
        ".pdf": ("pdf", parse_pdf),
        ".docx": ("docx", parse_docx),
        ".xlsx": ("xlsx", parse_xlsx),
        ".pptx": ("pptx", parse_pptx),
    }
    if extension not in parsers:
        raise ValueError(
            f"unsupported document extension: {extension or 'none'}"
        )
    parser_name, parser = parsers[extension]
    return {
        "parser": parser_name,
        "content": normalize_text(parser(file_path)),
    }


def main():
    if len(sys.argv) not in (2, 3):
        raise SystemExit("usage: parse_attachment.py <file-path> [extension]")
    extension = sys.argv[2] if len(sys.argv) == 3 else None
    print(
        json.dumps(parse_document(sys.argv[1], extension), ensure_ascii=False)
    )


if __name__ == "__main__":
    main()
